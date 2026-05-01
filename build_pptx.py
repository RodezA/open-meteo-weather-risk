"""
Build the Open-Meteo presentation as a .pptx file.
Run: python build_pptx.py

"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
import os

LOGO_PATH = os.path.join(os.path.dirname(__file__), "open_meteo_logo.png")

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x1a, 0x1a, 0x2e)
BLUE        = RGBColor(0x4a, 0x6f, 0xa5)
BLUE_LIGHT  = RGBColor(0xe8, 0xed, 0xf5)
GREEN       = RGBColor(0x2d, 0x7a, 0x2d)
GREEN_LIGHT = RGBColor(0xe6, 0xf4, 0xe6)
CAUTION     = RGBColor(0xb8, 0x7a, 0x00)
CAUTION_LT  = RGBColor(0xff, 0xf3, 0xcd)
STOP        = RGBColor(0xcc, 0x22, 0x22)
STOP_LIGHT  = RGBColor(0xfd, 0xe8, 0xe8)
WHITE       = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY  = RGBColor(0xf8, 0xf9, 0xfa)
MID_GRAY    = RGBColor(0x44, 0x44, 0x44)
BORDER      = RGBColor(0xdd, 0xe0, 0xec)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout


def add_slide(light_bg=False):
    sl = prs.slides.add_slide(BLANK)
    # Logo: bottom-right corner, always on a dark pill so white logo is visible
    lw, lh = Inches(1.7), Inches(0.5)
    lx = W - lw - Inches(0.25)
    ly = H - lh - Inches(0.15)
    if light_bg:
        pill = sl.shapes.add_shape(1, lx - Inches(0.1), ly - Inches(0.07),
                                   lw + Inches(0.2), lh + Inches(0.14))
        pill.fill.solid()
        pill.fill.fore_color.rgb = DARK_NAVY
        pill.line.fill.background()
    sl.shapes.add_picture(LOGO_PATH, lx, ly, lw, lh)
    return sl


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h,
        bg_color=None, border_color=None, border_width=Pt(1),
        radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = border_width
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=Pt(18), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def eyebrow(slide, text, light=False):
    c = RGBColor(0xaa, 0xbb, 0xdd) if not light else RGBColor(0x88, 0x88, 0xaa)
    txt(slide, text.upper(), 0.6, 0.3, 12, 0.35, size=Pt(10), bold=True, color=c)


def divider(slide, x, y, w=0.5, color=BLUE):
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def heading(slide, text, y=0.8, color=WHITE, size=Pt(40)):
    txt(slide, text, 0.6, y, 12, 1.4, size=size, bold=True, color=color)


def card(slide, x, y, w, h, title, body_lines,
         bg_col=WHITE, title_col=BLUE, body_col=MID_GRAY,
         border_col=BORDER, accent_col=None):
    # background
    s = box(slide, x, y, w, h, bg_color=bg_col, border_color=border_col)
    if accent_col:
        accent = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Pt(4), Inches(h)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_col
        accent.line.fill.background()
    # title
    txt(slide, title.upper(), x + 0.15, y + 0.12, w - 0.3, 0.28,
        size=Pt(8), bold=True, color=title_col)
    # body
    body_text = "\n".join(body_lines)
    txt(slide, body_text, x + 0.15, y + 0.38, w - 0.3, h - 0.5,
        size=Pt(10.5), color=body_col)


def dark_card(slide, x, y, w, h, title, body_lines):
    card(slide, x, y, w, h, title, body_lines,
         bg_col=RGBColor(0x28, 0x2a, 0x45),
         title_col=RGBColor(0x9a, 0xbf, 0xe8),
         body_col=RGBColor(0xcc, 0xd5, 0xe8),
         border_col=RGBColor(0x3a, 0x3c, 0x5a))


def badge(slide, x, y, text, style="green"):
    colors = {
        "green":   (GREEN_LIGHT, GREEN),
        "caution": (CAUTION_LT,  CAUTION),
        "stop":    (STOP_LIGHT,  STOP),
        "blue":    (BLUE_LIGHT,  BLUE),
    }
    bg_c, fg_c = colors[style]
    b = box(slide, x, y, 1.0, 0.28, bg_color=bg_c, border_color=bg_c)
    txt(slide, text, x + 0.05, y + 0.02, 0.9, 0.26,
        size=Pt(10), bold=True, color=fg_c, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_NAVY)

# accent bar left
bar = sl.shapes.add_shape(1, 0, 0, Inches(0.08), H)
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

eyebrow(sl, "Open-Meteo · Staff Data Engineer Challenge")
txt(sl, "Construction Site\nWeather Risk", 0.6, 0.9, 8, 2.2,
    size=Pt(52), bold=True, color=WHITE)
txt(sl, "Hourly work-stoppage risk assessment\ngrounded in OSHA thresholds",
    0.6, 3.3, 8, 0.9, size=Pt(18), color=RGBColor(0xaa, 0xbb, 0xdd))
divider(sl, 0.6, 4.4)

# stats row
for i, (num, label) in enumerate([("4", "Risk Factors"), ("4", "Activity Types"), ("7d", "Forecast Window"), ("0", "API Keys Needed")]):
    sx = 0.6 + i * 2.1
    txt(sl, num, sx, 4.65, 1.8, 0.8, size=Pt(36), bold=True,
        color=RGBColor(0x7a, 0x9f, 0xd4), align=PP_ALIGN.CENTER)
    txt(sl, label.upper(), sx, 5.4, 1.8, 0.4, size=Pt(9), bold=True,
        color=RGBColor(0x66, 0x77, 0x99), align=PP_ALIGN.CENTER)

txt(sl, "Open-Meteo API  ·  Pure Python Engine  ·  Streamlit UI",
    0.6, 6.8, 12, 0.4, size=Pt(11),
    color=RGBColor(0x55, 0x66, 0x88), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide(light_bg=True)
bg(sl, LIGHT_GRAY)

eyebrow(sl, "Overview", light=True)
heading(sl, "45-Minute Agenda", color=DARK_NAVY, size=Pt(36))

items = [
    ("01", "Problem Definition",     "5 min",  "Why construction sites? How did we go from\n'weather data' to a specific, high-value decision?"),
    ("02", "Architecture Deep-Dive", "15 min", "Data flow, scoring model, AWS production design,\nstack rationale, and error-handling strategy."),
    ("03", "Prototype Demo",         "10 min", "Live walkthrough of the risk assessment UI\nand example output for real locations."),
    ("04", "Scalability & Leadership","15 min", "Team structure, gotchas, roadmap, Q&A on\nmentorship and long-term ownership."),
]
for i, (num, title, time, desc) in enumerate(items):
    cx = 0.5 + i * 3.2
    b = box(sl, cx, 1.8, 3.0, 4.7, bg_color=WHITE, border_color=BORDER)
    txt(sl, num,   cx+0.2, 1.95, 0.6, 0.55, size=Pt(28), bold=True, color=BLUE)
    txt(sl, title, cx+0.2, 2.55, 2.6, 0.6,  size=Pt(13), bold=True, color=DARK_NAVY)
    txt(sl, time,  cx+0.2, 3.1,  2.6, 0.35, size=Pt(10), color=CAUTION, bold=True)
    txt(sl, desc,  cx+0.2, 3.5,  2.65, 2.6, size=Pt(10), color=MID_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM DEFINITION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide(light_bg=True)
bg(sl, LIGHT_GRAY)

eyebrow(sl, "Problem Definition · 01", light=True)
heading(sl, "Weather doesn't stop construction.\nUnmanaged weather risk does.", color=DARK_NAVY, size=Pt(30))

points = [
    ("⚠", "~30% of construction delays are weather-related",
          "yet most site supervisors check a generic forecast — not activity-specific thresholds"),
    ("⚖", "OSHA thresholds are precise, not ranges",
          "crane operations must halt at sustained wind >30 mph (1926.1417); electrical at any precipitation (1910.333)"),
    ("💡", "The insight gap",
          "a raw forecast says '25 mph wind.' It does not say 'suspend crane operations at 14:00.' That translation is the product."),
]
for i, (icon, title, sub) in enumerate(points):
    y = 2.5 + i * 1.35
    txt(sl, icon, 0.55, y, 0.5, 0.6, size=Pt(22), color=BLUE)
    txt(sl, title, 1.1, y, 7.0, 0.4, size=Pt(14), bold=True, color=DARK_NAVY)
    txt(sl, sub,   1.1, y+0.38, 7.0, 0.7, size=Pt(11), color=MID_GRAY)

# right callout
box(sl, 9.0, 2.4, 3.8, 4.5, bg_color=WHITE, border_color=BORDER)
box(sl, 9.0, 2.4, Inches(3.8), Pt(4),
    bg_color=BLUE, border_color=BLUE)
txt(sl, "THE QUESTION WE'RE\nANSWERING", 9.15, 2.55, 3.5, 0.55,
    size=Pt(9), bold=True, color=BLUE)
txt(sl, '"Which hours today can my crew actually work — and what specifically has to stop?"',
    9.15, 3.1, 3.5, 1.4, size=Pt(13), bold=True, color=DARK_NAVY)
txt(sl, "WHY I CHOSE CONSTRUCTION SITES", 9.15, 4.6, 3.5, 0.35,
    size=Pt(9), bold=True, color=BLUE)
txt(sl, "High-stakes, time-sensitive, OSHA-regulated. The thresholds are public, precise, and non-negotiable — ideal for a rules-based engine where correctness is verifiable.",
    9.15, 4.95, 3.5, 1.6, size=Pt(10), color=MID_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATA FLOW
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_NAVY)

eyebrow(sl, "Architecture Deep-Dive · 02")
heading(sl, "Prototype Data Flow", size=Pt(34))

boxes = [
    ("Open-Meteo API",   "free · no key\nhourly forecast"),
    ("Weather Client",   "async httpx\nPydantic models"),
    ("Risk Engine",      "pure Python\nzero deps"),
    ("Streamlit UI",     "chart + table\nactivity flags"),
]
bw = 2.4
gap = 0.3
start_x = (13.33 - (len(boxes)*bw + (len(boxes)-1)*gap)) / 2

for i, (title, sub) in enumerate(boxes):
    x = start_x + i*(bw + gap)
    highlight = i == 2
    bc = RGBColor(0x3a, 0x5a, 0x85) if highlight else RGBColor(0x28, 0x2a, 0x45)
    br = RGBColor(0x7a, 0x9f, 0xd4) if highlight else RGBColor(0x3a, 0x3c, 0x5a)
    b = box(sl, x/96, 2.5, bw/96*96, 1.5, bg_color=bc, border_color=br)

    # reuse inch-based coords
    bxi = x / 96
    txt(sl, title, bxi+0.1, 2.65, bw/96*96-0.2, 0.45,
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, sub, bxi+0.1, 3.1, bw/96*96-0.2, 0.6,
        size=Pt(10), color=RGBColor(0x99, 0xaa, 0xcc), align=PP_ALIGN.CENTER)

    if i < len(boxes)-1:
        ax = bxi + bw/96*96
        txt(sl, "→", ax+0.05, 3.0, 0.28, 0.4,
            size=Pt(20), color=RGBColor(0x55,0x66,0x88), align=PP_ALIGN.CENTER)

dark_card(sl, 0.5, 4.4, 3.8, 2.55,
    "KEY DESIGN CHOICE",
    ["The risk engine has zero framework",
     "dependencies. It takes HourlyWeather",
     "and returns List[HourlyRisk].",
     "",
     "Runs identically in Lambda, a cron job,",
     "or this UI — no refactor to scale."])

dark_card(sl, 4.6, 4.4, 3.8, 2.55,
    "COMPOSITE SCORE LOGIC",
    ["Four additive factors, each capped.",
     "Thunderstorm (WMO 95/96/99) alone",
     "contributes 70 pts — instant STOP.",
     "",
     "Score is capped at 100 to keep",
     "the scale intuitive and bounded."])

dark_card(sl, 8.7, 4.4, 3.8, 2.55,
    "ACTIVITY vs. COMPOSITE SCORE",
    ["The banner score drives the site-level",
     "risk level (GREEN / CAUTION / STOP).",
     "",
     "Activity restrictions use independent",
     "threshold comparisons — electrical",
     "doesn't wait for a CAUTION score."])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SCORING MODEL
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide(light_bg=True)
bg(sl, LIGHT_GRAY)

eyebrow(sl, "Architecture Deep-Dive · 02", light=True)
heading(sl, "Risk Scoring Model", color=DARK_NAVY, size=Pt(34))

# Left — factor table
box(sl, 0.4, 1.8, 6.0, 3.8, bg_color=WHITE, border_color=BORDER)
headers = ["Factor", "Max pts", "STOP trigger"]
col_x = [0.55, 2.9, 3.9]
for ci, (h, cx) in enumerate(zip(headers, col_x)):
    txt(sl, h.upper(), cx, 1.9, 1.3, 0.3,
        size=Pt(8), bold=True, color=RGBColor(0x88,0x88,0xaa), align=PP_ALIGN.LEFT)

rows = [
    ("Wind / Gusts",    "70", "≥45 mph effective"),
    ("Weather Code",    "70", "Thunderstorm (WMO 95/96/99)"),
    ("Precipitation",   "30", ">0.3 mm active"),
    ("Visibility",      "20", "<200 m"),
]
for ri, (f, pts, trigger) in enumerate(rows):
    ry = 2.3 + ri * 0.78
    if ri % 2 == 0:
        box(sl, 0.4, ry-0.05, 6.0, 0.75, bg_color=LIGHT_GRAY, border_color=LIGHT_GRAY)
    txt(sl, f,       0.55, ry, 2.2, 0.5, size=Pt(11), bold=True, color=DARK_NAVY)
    txt(sl, pts,     2.9,  ry, 0.8, 0.5, size=Pt(11), color=BLUE, align=PP_ALIGN.CENTER)
    txt(sl, trigger, 3.9,  ry, 2.3, 0.5, size=Pt(10), color=MID_GRAY)

# Right — activity table
box(sl, 6.8, 1.8, 6.0, 3.8, bg_color=WHITE, border_color=BORDER)
txt(sl, "ACTIVITY RESTRICTIONS (OSHA)", 6.95, 1.9, 5.6, 0.3,
    size=Pt(8), bold=True, color=RGBColor(0x88,0x88,0xaa))
act_rows = [
    ("🏗  Crane",      "Wind >30 mph or gusts >35 mph",   "OSHA 1926.1417"),
    ("🦺  Exterior",   "Wind >40 mph or thunderstorm",     ""),
    ("⚡  Electrical", "Any active precipitation",         "OSHA 1910.333"),
    ("👷  General",    "STOP-level conditions or thunderstorm", ""),
]
for ri, (act, cond, ref) in enumerate(act_rows):
    ry = 2.3 + ri * 0.78
    if ri % 2 == 0:
        box(sl, 6.8, ry-0.05, 6.0, 0.75, bg_color=LIGHT_GRAY, border_color=LIGHT_GRAY)
    txt(sl, act,  6.95, ry, 2.0, 0.5, size=Pt(11), bold=True, color=DARK_NAVY)
    txt(sl, cond, 9.0,  ry, 3.5, 0.35, size=Pt(10), color=MID_GRAY)
    if ref:
        txt(sl, ref, 9.0, ry+0.32, 3.5, 0.3, size=Pt(8), color=BLUE, italic=True)

# score bands bottom
for i, (label, score, style, bg_c, fg_c) in enumerate([
    ("GREEN",   "0–39",   "Normal operations",           GREEN_LIGHT, GREEN),
    ("CAUTION", "40–69",  "Elevated — supervisors alerted", CAUTION_LT, CAUTION),
    ("STOP",    "70–100", "Work suspension required",    STOP_LIGHT, STOP),
]):
    bx = 0.4 + i * 4.3
    b = box(sl, bx, 5.85, 4.0, 1.3, bg_color=bg_c, border_color=bg_c)
    txt(sl, label, bx+0.2, 5.95, 1.4, 0.45, size=Pt(18), bold=True, color=fg_c)
    txt(sl, score, bx+1.6, 6.0, 1.2, 0.35, size=Pt(11), color=fg_c, bold=True, align=PP_ALIGN.RIGHT)
    txt(sl, style, bx+0.2, 6.45, 3.6, 0.4, size=Pt(10), color=fg_c)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TECH STACK RATIONALE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_NAVY)

eyebrow(sl, "Architecture Deep-Dive · 02")
heading(sl, "Tech Stack Rationale", size=Pt(34))

choices = [
    ("Pure Python Engine",
     "Why not Pandas/Spark?",
     ["Scoring logic is O(n) per hour — no\n",
      "vectorization needed for <168 rows.\n",
      "Pure dataclasses = zero import overhead\n",
      "in Lambda cold start."]),
    ("Async httpx",
     "Why not requests?",
     ["Open-Meteo calls are I/O-bound. Async\n",
      "allows future concurrency (batch fetch\n",
      "multiple locations) with no refactor.\n",
      "Sync wrapper via asyncio.run() for Streamlit."]),
    ("Streamlit (Prototype)",
     "Why not React/FastAPI?",
     ["Fastest path to a demo that a non-engineer\n",
      "can run with one command. Engine is\n",
      "framework-agnostic — UI is swappable\n",
      "without touching business logic."]),
    ("Pydantic Models",
     "Why strict typing?",
     ["Open-Meteo returns raw JSON. Pydantic\n",
      "validates at the boundary — if the API\n",
      "changes shape, we get an explicit error\n",
      "instead of a silent wrong score."]),
]
for i, (title, why, body) in enumerate(choices):
    cx = 0.4 + i * 3.15
    dark_card(sl, cx, 2.2, 2.9, 4.7, why.upper(), [title] + [""] + body)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ERROR HANDLING & OPERATIONAL CONCERNS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide(light_bg=True)
bg(sl, LIGHT_GRAY)

eyebrow(sl, "Architecture Deep-Dive · 02", light=True)
heading(sl, "Error Handling &\nOperational Concerns", color=DARK_NAVY, size=Pt(32))

items = [
    (STOP,    "API Failure",
     "httpx raises on non-2xx. Streamlit catches and\ndisplays a user-facing error — never a silent wrong\nresult. In production: DLQ replay, SNS alert to on-call."),
    (CAUTION, "Open-Meteo Rate Limits",
     "Free tier is ~10k calls/day. At 15-min refresh × 10k\nsites = 960 calls/hour — well within limits. SQS\nbackpressure prevents burst beyond quota."),
    (BLUE,    "Stale Cache Risk",
     "DynamoDB TTL set to 2× the refresh interval. If the\ningest Lambda is down, API returns the last known score\nwith a staleness timestamp — never silent null data."),
    (GREEN,   "Threshold Correctness",
     "14 unit tests lock every scoring boundary. OSHA\nthresholds are stored in SSM (production) so safety\nmanagers can update without a code deployment."),
]
for i, (color, title, body) in enumerate(items):
    cx = 0.5 + i * 3.15
    b = box(sl, cx, 2.4, 2.9, 4.5, bg_color=WHITE, border_color=BORDER)
    accent = sl.shapes.add_shape(1, Inches(cx), Inches(2.4), Pt(5), Inches(4.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = color; accent.line.fill.background()
    txt(sl, title, cx+0.18, 2.52, 2.55, 0.45, size=Pt(13), bold=True, color=DARK_NAVY)
    txt(sl, body,  cx+0.18, 3.05, 2.6, 3.5, size=Pt(10.5), color=MID_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ACTIONABLE OUTPUT
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_NAVY)

eyebrow(sl, "Architecture Deep-Dive · 02")
heading(sl, "Actionable vs. Reformatted", size=Pt(34))

# Left — raw forecast
box(sl, 0.4, 2.0, 5.8, 4.8, bg_color=RGBColor(0x22,0x24,0x3a), border_color=RGBColor(0x55,0x22,0x22))
txt(sl, "❌  RAW FORECAST", 0.6, 2.1, 5.4, 0.4, size=Pt(10), bold=True, color=STOP)
raw_lines = [
    "Monday 14:00",
    "Temperature: 72°F",
    "Wind: 28 mph",
    "Gusts: 38 mph",
    "Precipitation: 0 mm",
    "Visibility: 8 miles",
    "Conditions: Partly cloudy",
]
for i, line in enumerate(raw_lines):
    txt(sl, line, 0.65, 2.62+i*0.47, 5.3, 0.42, size=Pt(12), color=RGBColor(0x99,0xaa,0xbb))

txt(sl, "A supervisor reads this and still has to\nanswer: 'Can I run the crane? Do I need\nto stop electrical?' The data is there.\nThe decision is not.",
    0.65, 5.75, 5.3, 1.2, size=Pt(10), color=RGBColor(0x77,0x88,0x99), italic=True)

# Right — system output
box(sl, 6.9, 2.0, 5.8, 4.8, bg_color=RGBColor(0x1a,0x2a,0x1a), border_color=RGBColor(0x22,0x55,0x22))
txt(sl, "✓  OPEN-METEO OUTPUT", 7.1, 2.1, 5.4, 0.4, size=Pt(10), bold=True, color=GREEN)
txt(sl, "Monday 14:00", 7.15, 2.6, 5.3, 0.4, size=Pt(12), bold=True, color=WHITE)

b = box(sl, 7.15, 3.08, 2.3, 0.45, bg_color=CAUTION_LT, border_color=CAUTION_LT)
txt(sl, "CAUTION  ·  Score 48", 7.2, 3.12, 2.2, 0.35, size=Pt(11), bold=True, color=CAUTION, align=PP_ALIGN.CENTER)

txt(sl, "Primary driver: Gusts 38 mph / sustained 28 mph", 7.15, 3.65, 5.3, 0.4, size=Pt(10.5), color=RGBColor(0x99,0xcc,0x99))

acts = [("🏗 Crane", "SUSPENDED — gusts 38 mph exceed OSHA 1926.1417", STOP, STOP_LIGHT),
        ("🦺 Exterior", "Allowed — PPE check required", CAUTION, CAUTION_LT),
        ("⚡ Electrical", "Allowed — no active precipitation", GREEN, GREEN_LIGHT),
        ("👷 General", "Allowed — score below STOP threshold", GREEN, GREEN_LIGHT)]
for i, (act, reason, fc, bc) in enumerate(acts):
    ry = 4.15 + i*0.47
    b2 = box(sl, 7.15, ry, 5.3, 0.4, bg_color=bc, border_color=bc)
    txt(sl, act,    7.25, ry+0.04, 1.2, 0.34, size=Pt(9),  bold=True, color=fc)
    txt(sl, reason, 8.5,  ry+0.04, 3.8, 0.34, size=Pt(9),  color=fc)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AWS ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide(light_bg=True)
bg(sl, LIGHT_GRAY)

eyebrow(sl, "Architecture Deep-Dive · 02", light=True)
heading(sl, "AWS Production Architecture", color=DARK_NAVY, size=Pt(32))

# Two columns: ingest path + request path
box(sl, 0.35, 2.0, 6.1, 5.1, bg_color=WHITE, border_color=BORDER)
box(sl, 0.35, 2.0, Inches(6.1), Pt(4), bg_color=BLUE, border_color=BLUE)
txt(sl, "INGEST PATH  (scheduled every 15 min)", 0.55, 2.08, 5.7, 0.35,
    size=Pt(9), bold=True, color=BLUE)

ingest_steps = [
    ("EventBridge Scheduler", "Fires every 15 min for all registered sites"),
    ("SQS FIFO Queue",        "Queues location refresh jobs with backpressure"),
    ("Lambda — ingest",       "Fetches Open-Meteo, runs risk engine, writes scores"),
    ("DynamoDB",              "Caches risk by (lat, lon, hour) with TTL"),
    ("S3 Parquet Archive",    "Raw weather for replay, audit trail, ML data"),
    ("Lambda — alert",        "DynamoDB Streams → SNS → email/Slack on STOP/CAUTION"),
]
for i, (comp, desc) in enumerate(ingest_steps):
    ry = 2.55 + i * 0.72
    if i < len(ingest_steps)-1:
        txt(sl, "↓", 0.75, ry+0.38, 0.3, 0.28, size=Pt(11), color=RGBColor(0xaa,0xaa,0xcc), align=PP_ALIGN.CENTER)
    txt(sl, comp, 1.1, ry, 2.5, 0.38, size=Pt(11), bold=True, color=DARK_NAVY)
    txt(sl, desc, 3.7, ry, 2.5, 0.38, size=Pt(10), color=MID_GRAY)

box(sl, 6.75, 2.0, 6.1, 5.1, bg_color=WHITE, border_color=BORDER)
box(sl, 6.75, 2.0, Inches(6.1), Pt(4), bg_color=GREEN, border_color=GREEN)
txt(sl, "REQUEST PATH  (on-demand API)", 6.95, 2.08, 5.7, 0.35,
    size=Pt(9), bold=True, color=GREEN)

req_steps = [
    ("CloudFront + API Gateway", "Edge caching + REST endpoint"),
    ("Lambda — risk",            "Reads DynamoDB cache → sub-10ms response"),
    ("Cache miss",               "Enqueues SQS job for background backfill"),
    ("SSM Parameter Store",      "OSHA thresholds — no code deploy to update"),
    ("CloudWatch + X-Ray",       "Logs, metrics, distributed tracing"),
    ("Dead Letter Queue",        "Failed location refreshes captured for replay"),
]
for i, (comp, desc) in enumerate(req_steps):
    ry = 2.55 + i * 0.72
    if i < len(req_steps)-1:
        txt(sl, "↓", 7.15, ry+0.38, 0.3, 0.28, size=Pt(11), color=RGBColor(0xaa,0xaa,0xcc), align=PP_ALIGN.CENTER)
    txt(sl, comp, 7.5, ry, 2.7, 0.38, size=Pt(11), bold=True, color=DARK_NAVY)
    txt(sl, desc, 10.25, ry, 2.3, 0.38, size=Pt(10), color=MID_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DEMO SLIDE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_NAVY)

eyebrow(sl, "Prototype Demo · 03")
heading(sl, "Live Demo", size=Pt(40))

txt(sl, "streamlit run streamlit_app.py", 0.6, 2.2, 8, 0.55,
    size=Pt(18), color=RGBColor(0x7a, 0x9f, 0xd4))

dark_card(sl, 0.4, 3.1, 3.9, 3.7,
    "WHAT TO SHOW",
    ["1. New York — calm conditions (GREEN)",
     "2. Miami in summer — heat stress",
     "3. Chicago — wind thresholds, crane",
     "   suspension at CAUTION boundary",
     "4. 7-day forecast — peak risk banner",
     "5. OSHA reference expander"])

dark_card(sl, 4.6, 3.1, 3.9, 3.7,
    "TALKING POINTS",
    ["Each hour shows WHY, not just WHAT",
     "Activity restrictions fire independently",
     "Primary driver is explicit — not inferred",
     "Score is additive — show two moderate",
     "  factors combining to CAUTION",
     "Zero API keys — runs immediately"])

dark_card(sl, 8.8, 3.1, 3.9, 3.7,
    "ANTICIPATED QUESTIONS",
    ["Q: What if Open-Meteo is down?",
     "→ Cached results + DLQ replay",
     "",
     "Q: Can thresholds be customized?",
     "→ SSM in production; per-site overrides",
     "  are on the roadmap"])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TEAM OF 3
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide(light_bg=True)
bg(sl, LIGHT_GRAY)

eyebrow(sl, "Scalability & Leadership · 04", light=True)
heading(sl, "Leading a Team of 3 to Production", color=DARK_NAVY, size=Pt(30))

teams = [
    ("Engineer 1\nData Infra",
     ["SQS + EventBridge ingest pipeline",
      "DynamoDB schema + TTL tuning",
      "S3 Parquet archive + Glue catalog",
      "DLQ retry logic",
      "Load & performance testing"]),
    ("Engineer 2\nAPI + Engine",
     ["Lambda packaging + IaC (CDK)",
      "SSM threshold config layer",
      "API Gateway + Cognito auth",
      "Boundary & integration tests",
      "Multi-source weather fallback"]),
    ("Engineer 3\nUI + Alerts",
     ["Geocoding (address → coordinates)",
      "SNS/SES/Slack alert subscriptions",
      "Historical risk replay from S3",
      "Mobile-responsive UI",
      "Per-site threshold override UI"]),
]
for i, (title, items) in enumerate(teams):
    cx = 0.4 + i * 4.15
    box(sl, cx, 2.0, 3.85, 4.85, bg_color=WHITE, border_color=BORDER)
    box(sl, cx, 2.0, Inches(3.85), Pt(4), bg_color=BLUE, border_color=BLUE)
    txt(sl, title, cx+0.15, 2.1, 3.5, 0.7, size=Pt(14), bold=True, color=DARK_NAVY)
    for j, item in enumerate(items):
        txt(sl, "·  " + item, cx+0.15, 2.92+j*0.72, 3.5, 0.6, size=Pt(11), color=MID_GRAY)

txt(sl, "Week 1–2: IaC skeleton + DynamoDB schema (Eng 1) in parallel with Lambda packaging (Eng 2)  ·  Week 3–4: Ingest pipeline live, API serving cached results  ·  Week 5–6: Alerting + UI polish + load test",
    0.4, 7.1, 12.5, 0.4, size=Pt(9), color=RGBColor(0x88,0x88,0xaa), italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — GOTCHAS + ROADMAP
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide(light_bg=True)
bg(sl, LIGHT_GRAY)

eyebrow(sl, "Scalability & Leadership · 04", light=True)
heading(sl, "Gotchas & Roadmap", color=DARK_NAVY, size=Pt(34))

box(sl, 0.4, 2.05, 6.05, 4.95, bg_color=WHITE, border_color=BORDER)
box(sl, 0.4, 2.05, Inches(6.05), Pt(4), bg_color=STOP, border_color=STOP)
txt(sl, "GOTCHAS TO WARN THE TEAM", 0.6, 2.13, 5.6, 0.35,
    size=Pt(9), bold=True, color=STOP)
gotchas = [
    ("Open-Meteo updates hourly",
     "Polling faster than 15 min wastes quota and returns identical data."),
    ("OSHA limits are minimums",
     "Crane manufacturer limits can be lower — support per-site threshold overrides."),
    ("Effective wind formula",
     "max(sustained, gusts × 0.8) is a judgment call. Document for safety review."),
    ("DynamoDB hot partitions",
     "10k sites refreshing at :00 every hour → jitter the EventBridge schedule."),
    ("Thunderstorm code coverage",
     "WMO codes 95/96/99 only. Drizzle, freezing rain, fog not yet scored."),
]
for i, (title, sub) in enumerate(gotchas):
    ry = 2.6 + i * 0.84
    txt(sl, "!", 0.6, ry, 0.3, 0.4, size=Pt(13), bold=True, color=STOP)
    txt(sl, title, 0.95, ry, 5.0, 0.35, size=Pt(11), bold=True, color=DARK_NAVY)
    txt(sl, sub, 0.95, ry+0.36, 5.1, 0.38, size=Pt(10), color=MID_GRAY)

box(sl, 6.75, 2.05, 6.05, 4.95, bg_color=WHITE, border_color=BORDER)
box(sl, 6.75, 2.05, Inches(6.05), Pt(4), bg_color=BLUE, border_color=BLUE)
txt(sl, "ROADMAP — WHAT I'D BUILD NEXT", 6.95, 2.13, 5.6, 0.35,
    size=Pt(9), bold=True, color=BLUE)
roadmap = [
    ("Sprint 1",   "SSM threshold config · multi-source weather fallback (NWS API) · geocoding"),
    ("Sprint 2",   "Per-site threshold overrides · historical replay from S3 archive"),
    ("Sprint 3",   "Alert subscriptions (email/Slack) · mobile UI · batch site import (CSV)"),
    ("Later",      "ML risk model trained on S3 archive + incident reports · temp/humidity scoring"),
    ("Prototype\ntrade-offs",  "FastAPI removed for simplicity · no auth · thresholds hardcoded"),
]
for i, (phase, items) in enumerate(roadmap):
    ry = 2.6 + i * 0.84
    txt(sl, phase, 6.95, ry, 1.3, 0.7, size=Pt(11), bold=True, color=BLUE)
    txt(sl, items, 8.3, ry, 4.3, 0.7, size=Pt(10), color=MID_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CLOSE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, BLUE)

bar2 = sl.shapes.add_shape(1, 0, 0, Inches(0.08), H)
bar2.fill.solid(); bar2.fill.fore_color.rgb = DARK_NAVY; bar2.line.fill.background()

eyebrow(sl, "Summary")
txt(sl, "From raw forecast\nto actionable decision.", 0.6, 1.0, 11, 2.4,
    size=Pt(48), bold=True, color=WHITE)
divider(sl, 0.6, 3.6, color=WHITE)

summary = [
    ("Engine",     "Pure Python · zero deps · independently testable · Lambda-ready"),
    ("Insight",    "Not just a score — a specific suspension with an OSHA citation"),
    ("Scale Path", "EventBridge → SQS → Lambda → DynamoDB · 10k sites, no refactor"),
    ("Tests",      "14 unit tests lock every scoring boundary and activity threshold"),
]
for i, (label, body) in enumerate(summary):
    cx = 0.5 + i * 3.1
    b = box(sl, cx, 4.0, 2.85, 2.7,
            bg_color=RGBColor(0x3a,0x5a,0x85),
            border_color=RGBColor(0x5a,0x7a,0xa5))
    txt(sl, label, cx+0.15, 4.1, 2.55, 0.4, size=Pt(10), bold=True,
        color=RGBColor(0xcc,0xdd,0xf5))
    txt(sl, body, cx+0.15, 4.55, 2.55, 1.8, size=Pt(10.5), color=WHITE)

txt(sl, "Open-Meteo API  ·  Python  ·  Streamlit  ·  AWS (future state)",
    0.6, 7.1, 12.1, 0.4, size=Pt(10),
    color=RGBColor(0xaa,0xbb,0xdd), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
prs.save("open_meteo_presentation.pptx")
print("Saved: open_meteo_presentation.pptx")
